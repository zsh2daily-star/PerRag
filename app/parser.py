"""文档解析模块 —— 支持 PDF、Word、Excel、PPT、Markdown、纯文本等多种格式。

每种文件类型都有对应的解析函数，统一的入口是 parse_file()。
解析结果统一封装为 ParsedDocument 对象，供 indexer.py 使用。
"""

from __future__ import annotations

import logging
import mimetypes
import re
from dataclasses import dataclass, field
from pathlib import Path

import httpx
from pypdf import PdfReader

from app.config import settings
from app.mineru_client import MinerUError, parse_pdf_with_mineru

logger = logging.getLogger(__name__)


# 支持的所有文件扩展名
SUPPORTED_EXTENSIONS = {".pdf", ".md", ".markdown", ".txt", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"}

# 文本文件编码探测顺序（中文文件常见编码，按使用频率排列）
_ENCODING_CANDIDATES = (
    "utf-8",       # 最通用
    "gb18030",     # 中文国标（向下兼容 GBK、GB2312）
    "gbk",
    "gb2312",
    "big5",        # 繁体中文
    "shift_jis",   # 日文
    "cp1252",      # Windows 西欧编码
    "latin-1",     # 兜底：几乎都能解码，但有乱码风险
)


# ── 数据结构 ──────────────────────────────────────────────


@dataclass
class ParsedDocument:
    """解析后的文档 —— 索引模块使用的标准格式。

    所有解析函数都返回这个对象，保证接口统一。
    """

    text: str                        # 文档全文（已清洗）
    metadata: dict = field(default_factory=dict)   # 文档属性（标题、作者等）
    page_count: int | None = None                  # 页数（PDF 专有，其他为 None）
    source_format: str = "text"                    # 来源格式标记（pdf/docx/markdown...）

    @property
    def is_empty(self) -> bool:
        """检查文档是否为空。"""
        return not self.text.strip()


# ── 编码检测 ──────────────────────────────────────────────


def _detect_encoding(path: Path) -> str:
    """自动检测文本文件的字符编码。

    策略:
    1. 优先用 chardet 库（如果已安装）—— 准确率最高
    2. 降级：遍历常见中文编码列表逐一尝试解码
    3. 都失败则返回 utf-8（UnicodeDecodeError 时用 errors="replace" 兜底）

    为什么需要这个？Windows 上保存的 .txt/.md 文件常是 GBK 编码，
    直接用 utf-8 读取会乱码。
    """
    raw = path.read_bytes()
    if not raw:
        return "utf-8"

    # 优先 chardet
    try:
        import chardet

        result = chardet.detect(raw)
        enc = result.get("encoding")
        if enc and result.get("confidence", 0) > 0.5:
            enc = enc.lower().replace("-", "")
            alias_map = {
                "gb2312": "gb18030",
                "gbk": "gb18030",
                "ascii": "utf-8",  # chardet 有时把纯英文识别为 ascii
            }
            return alias_map.get(enc, enc)
    except ImportError:
        pass

    # 降级：尝试已知编码列表
    for enc in _ENCODING_CANDIDATES:
        try:
            raw.decode(enc)
            return enc
        except (UnicodeDecodeError, LookupError):
            continue

    return "utf-8"


# ── 文件类型检测 ──────────────────────────────────────────


def detect_mime_type(path: Path) -> str | None:
    """通过扩展名检测文件的 MIME 类型。

    返回如 "application/pdf"、"application/vnd.openxmlformats-officedocument.wordprocessingml.document"。
    """
    mime, _ = mimetypes.guess_type(str(path))
    return mime


def _warn_extension_mismatch(path: Path, expected_category: str) -> None:
    """警告：文件扩展名与实际 MIME 类型不匹配。

    比如用户把 .docx 改名为 .txt 骗过扩展名检测，
    但 MIME 类型仍然是 Word 文档，此时给出警告。
    """
    mime = detect_mime_type(path)
    suffix = path.suffix.lower()
    if mime and expected_category not in mime:
        logger.warning(
            "文件扩展名 %s 与检测到的 MIME 类型 %s 不一致: %s",
            suffix, mime, path.name,
        )


# ── 解析质量检测 ──────────────────────────────────────────


def _sample_positions(text: str, sample_size: int = 1000, positions: int = 5) -> list[str]:
    """从文本的前、中、后等多段均匀取样，避免单点误判。

    为什么需要多段采样？古籍 PDF 可能前半本是文字（解析正常）、后半本是表格
    或图片（解析为空），或者反过来。多段采样能覆盖整篇文档。
    """
    total = len(text)
    if total <= sample_size * positions:
        return [text]

    samples = []
    for i in range(positions):
        start = int(total * i / positions)
        end = min(start + sample_size, total)
        samples.append(text[start:end])
    return samples


def check_parse_quality(text: str, sample_size: int = 1000) -> dict:
    """检测 MinerU 解析输出的文本质量。返回质量评分和诊断信息。

    核心思路：OCR 乱码 = 零星的非中文字符散布在中文字符之间，
    而不是正常的成片数字/英文。关键区分：

      OCR 乱码：   "取其經 il. 1 + tnt TL 4 11 .. 2 沁中"
                   ↑ 数字和拉丁字母是 1-2 个孤立散布的

      正常文档：   "产能 63.5GWh，eVTOL 低空经济 3000P"
                   ↑ 数字和英文是成片出现的，代表数据/缩写

    检测三个维度：
    - 孤立数字密度：两边被 CJK/空格包围的 1-2 位数字碎片
    - 孤立拉丁密度：两边被 CJK/空格包围的 1-2 个字母碎片
    - 行内空格密度：CJK 字符之间的无意义空格

    返回:
        {"score": 0.0~1.0, "is_garbled": bool, "details": str}
        score ≥ 0.5 认为正常，< 0.5 认为乱码。
    """

    if not text or len(text) < 100:
        return {"score": 0.0, "is_garbled": True, "details": "文本过短或为空"}

    def _is_cjk(ch: str) -> bool:
        return '一' <= ch <= '鿿' or '㐀' <= ch <= '䶿'

    samples = _sample_positions(text, sample_size)
    all_scores: list[float] = []

    for sample in samples:
        total = len(sample)

        cjk = 0
        isolated_digits = 0   # OCR 特征：零星散布的数字
        isolated_latin = 0    # OCR 特征：零星散布的拉丁字母
        inline_spaces = 0     # CJK 文本中的行内空格

        i = 0
        while i < total:
            ch = sample[i]
            if _is_cjk(ch):
                cjk += 1
                i += 1
                continue

            if ch.isdigit() or (ch.isascii() and ch.isalpha()):
                # 收集连续的同类字符序列
                seq_end = i
                is_digit_run = ch.isdigit()
                while seq_end < total and (
                    sample[seq_end].isdigit() if is_digit_run
                    else (sample[seq_end].isascii() and sample[seq_end].isalpha())
                ):
                    seq_end += 1

                seq_len = seq_end - i

                if is_digit_run and seq_len <= 2:
                    # 1-2 位短数字序列 → 检查是否孤立（两边是 CJK 或非法定界符）
                    left = sample[i-1] if i > 0 else ' '
                    right = sample[seq_end] if seq_end < total else ' '
                    if (_is_cjk(left) or left.isspace() or not left.isalnum()) and \
                       (_is_cjk(right) or right.isspace() or not right.isalnum()):
                        isolated_digits += seq_len
                elif not is_digit_run and seq_len <= 2:
                    # 1-2 字母短序列 → 检查是否孤立
                    left = sample[i-1] if i > 0 else ' '
                    right = sample[seq_end] if seq_end < total else ' '
                    if (_is_cjk(left) or left.isspace()) and \
                       (_is_cjk(right) or right.isspace()):
                        isolated_latin += seq_len

                i = seq_end
                continue

            if ch == ' ':
                # 只算 CJK 文本中的行内空格
                left_cjk = i > 0 and _is_cjk(sample[i-1])
                right_cjk = i + 1 < total and _is_cjk(sample[i+1])
                if left_cjk or right_cjk:
                    inline_spaces += 1

            i += 1

        # 计算比率
        iso_digit_ratio = isolated_digits / total
        iso_latin_ratio = isolated_latin / total
        space_ratio = inline_spaces / total
        cjk_ratio = cjk / total

        # 噪声评分：三个维度独立计分
        # 孤立数字 >3%   → 强乱码信号
        # 孤立拉丁 >2%   → 极强乱码信号（正常中文文档几乎不会出现）
        # 行内空格 >6%   → 中等乱码信号
        noise = (
            min(iso_digit_ratio / 0.03, 1.0) * 20  # 孤立数字超 3% 满分
            + min(iso_latin_ratio / 0.02, 1.0) * 30  # 孤立拉丁超 2% 满分
            + min(space_ratio / 0.06, 1.0) * 10     # 行内空格超 6% 满分
        )
        # 总分 0-60，>25 判乱码

        score = max(0.0, min(1.0, 1.0 - noise / 60.0))
        all_scores.append(score)

    # 取 P40（5 样本中倒数第二差）而非 min()：
    # - 正常文档偶尔有统计表/索引页导致单个样本分低
    # - 乱码文档整篇都差，5 个样本里至少 4 个低分
    # P40 能容忍单个异常段，不被统计表/英文摘要误触发
    sorted_scores = sorted(all_scores)
    p40_idx = max(0, int(len(sorted_scores) * 0.4))
    final_score = sorted_scores[p40_idx]

    # 额外保护：如果只有1个样本不及格（<0.5）但大多数及格，不判乱码
    bad_count = sum(1 for s in all_scores if s < 0.5)
    total_samples = len(all_scores)

    return {
        "score": round(final_score, 3),
        "is_garbled": final_score < 0.5 and bad_count >= total_samples * 0.4,
        "details": (
            f"cjk={cjk_ratio*100:.0f}% iso_digit={iso_digit_ratio*100:.1f}% "
            f"iso_latin={iso_latin_ratio*100:.1f}% sp={space_ratio*100:.0f}% "
            f"({bad_count}/{total_samples} bad)"
        ),
    }


# ── PDF 解析 ──────────────────────────────────────────────


def parse_pdf_with_pypdf(path: Path) -> ParsedDocument:
    """使用 pypdf 提取 PDF 文本（纯文字型 PDF）。

    局限性: 只能处理文字型 PDF，无法处理扫描件/图片型 PDF。
    对于扫描件，会抛出 ValueError，上游应切换 MinerU。
    """
    reader = PdfReader(str(path))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text.strip())

    if not pages:
        raise ValueError(f"pypdf 未提取到文本（可能是扫描件）: {path.name}")

    # 提取 PDF 文档属性（标题、作者等）
    metadata: dict[str, str] = {}
    info = reader.metadata
    if info:
        for key in ("/Title", "/Author", "/Subject", "/Creator"):
            value = getattr(info, key, None) or (info.get(key) if isinstance(info, dict) else None)
            if value:
                metadata[key.lstrip("/").lower()] = str(value).strip()

    return ParsedDocument(
        text="\n\n".join(pages),
        metadata=metadata,
        page_count=len(reader.pages),
        source_format="pdf",
    )


def parse_pdf(path: Path) -> ParsedDocument:
    """PDF 解析统一入口 —— 根据配置选择解析策略。

    三种模式（由 PDF_PARSER 环境变量控制）：
    - mineru : 只用 MinerU（GPU OCR，适合扫描件）
    - pypdf  : 只用 pypdf（快，但只能处理文字型 PDF）
    - auto   : 先 MinerU，失败后自动降级 pypdf

    MinerU 解析后会跑质量检测，不达标的在 metadata 里标记 quality_warning，
    不阻断索引流程——搜索时 LLM 能看到这个标记并告知用户结果可能不可靠。
    """
    _warn_extension_mismatch(path, "pdf")
    mode = settings.pdf_parser.lower()

    if mode == "pypdf":
        return parse_pdf_with_pypdf(path)

    if mode in {"mineru", "auto"}:
        try:
            text = parse_pdf_with_mineru(path)

            quality = check_parse_quality(text)
            logger.info(
                "MinerU 解析完成 %s: score=%.3f garbled=%s %s",
                path.name, quality["score"], quality["is_garbled"], quality["details"],
            )

            metadata: dict = {}
            if quality["is_garbled"]:
                metadata["quality_warning"] = (
                    f"本文档 OCR 解析质量低 (score={quality['score']:.2f})，"
                    f"可能存在大量乱码或缺失文字，内容仅供参考。"
                )
                logger.warning("⚠ OCR 疑似乱码: %s (score=%.2f)", path.name, quality["score"])

            return ParsedDocument(
                text=text,
                metadata=metadata,
                page_count=None,
                source_format="pdf",
            )
        except (MinerUError, httpx.HTTPError) as exc:
            if mode == "mineru":
                raise ValueError(f"MinerU 解析失败: {exc}") from exc
            logger.warning("MinerU 失败，降级 pypdf: %s -> %s", path.name, exc)

    return parse_pdf_with_pypdf(path)


# ── Word (.doc) 解析（旧格式）───────────────────────────────


def parse_doc(path: Path) -> ParsedDocument:
    """解析旧版 Word (.doc) 文档，使用 antiword 提取文本。

    antiword 是 Linux 下最常用的 .doc 转文本工具，无需 GUI。
    对 WPS 等非标准生成的 .doc，增加降级策略。
    """
    import subprocess

    # 策略 1：antiword（标准 .doc）
    text = _try_antiword(path)
    if text:
        return ParsedDocument(text=text, metadata={}, source_format="doc")

    # 策略 2：python-docx（.doc 实际是 .docx 但扩展名错了）
    text = _try_docx(path)
    if text:
        return ParsedDocument(text=text, metadata={}, source_format="doc")

    # 策略 3：olefile 读取 OLE2 WordDocument 流（WPS 等非标准 .doc）
    text = _try_olefile(path)
    if text:
        return ParsedDocument(text=text, metadata={}, source_format="doc")

    raise ValueError(f"无法解析 .doc 文件（antiword/docx/olefile 均失败）: {path.name}")


def _try_antiword(path: Path) -> str:
    import subprocess
    try:
        result = subprocess.run(
            ["antiword", "-m", "UTF-8.txt", str(path)],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except Exception:
        pass
    return ""


def _try_docx(path: Path) -> str:
    """尝试用 python-docx 解析（有些 .doc 实际是 .docx）。"""
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(str(path))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))
        return "\n".join(parts) if parts else ""
    except Exception:
        return ""


def _try_olefile(path: Path) -> str:
    """尝试用 olefile 读取 OLE2 WordDocument 流（WPS 等非标准 .doc）。"""
    try:
        from olefile import OleFileIO
        ole = OleFileIO(str(path))
        if ole.exists("WordDocument"):
            data = ole.openstream("WordDocument").read()
            # Word FIB (File Information Block) 中 text 起始偏移在 0x18 处
            # 这里简单取可打印字符段
            import re
            text = data.decode("utf-16-le", errors="ignore")
            # 过滤出中英文可读内容
            segments = re.findall(r"[一-鿿　-〿＀-￯a-zA-Z0-9\s　-〿＀-￯\-.,()\[\]{}:;!?@#$%^&*+=/\\|<>`~]{4,}", text)
            result = "\n".join(segments)
            return result if len(result) > 100 else ""
    except Exception:
        return ""


# ── Word (.docx) 解析 ─────────────────────────────────────


def parse_docx(path: Path) -> ParsedDocument:
    """解析 Word (.docx) 文档，提取段落与表格中的文本。

    .docx 本质是一个 ZIP 压缩包，内含 XML 文件。
    python-docx 库负责解析这些 XML，提取出文字内容。

    注意: .doc（旧格式）请使用 parse_doc()，本函数仅处理 .docx。
    """
    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise ImportError("解析 .docx 需要 python-docx 库") from None

    doc = DocxDocument(str(path))
    parts: list[str] = []

    # 提取所有段落文本
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 提取所有表格内容（用 | 分隔单元格）
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))

    if not parts:
        raise ValueError(f"无法从 .docx 文件中提取文本: {path.name}")

    # 提取文档属性
    metadata: dict[str, str] = {}
    props = doc.core_properties
    for attr in ("title", "author", "subject", "last_modified_by"):
        value = getattr(props, attr, None)
        if value:
            metadata[attr] = str(value).strip()

    return ParsedDocument(
        text="\n\n".join(parts),
        metadata=metadata,
        source_format="docx",
    )


# ── Excel (.xlsx / .xls) 解析 ─────────────────────────────


def parse_xlsx(path: Path) -> ParsedDocument:
    """解析 Excel 文件，提取所有 sheet 中的文本。

    每个 sheet 以 "--- sheet名 ---" 开头，
    每行数据用 " | " 分隔。
    忽略空单元格，跳过完全空的行。
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("解析 .xlsx 需要 openpyxl 库") from None

    # read_only=True: 只读模式，内存占用更低
    # data_only=True: 读取公式的计算结果而非公式本身
    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    sheet_names: list[str] = []

    for sheet in wb.worksheets:
        sheet_names.append(sheet.title)
        sheet_parts: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                sheet_parts.append(" | ".join(cells))
        if sheet_parts:
            parts.append(f"--- {sheet.title} ---")
            parts.extend(sheet_parts)

    wb.close()

    if not parts:
        raise ValueError(f"无法从 Excel 文件中提取文本: {path.name}")

    return ParsedDocument(
        text="\n\n".join(parts),
        metadata={"sheets": ", ".join(sheet_names)},
        source_format="xlsx",
    )


# ── PowerPoint (.pptx /.ppt) 解析 ─────────────────────────


def parse_pptx(path: Path) -> ParsedDocument:
    """解析 PowerPoint 文件，提取每张幻灯片中的文本。

    每张幻灯片以 "--- Slide N ---" 分隔。
    同时提取幻灯片中的表格（如有）。
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("解析 .pptx 需要 python-pptx 库") from None

    prs = Presentation(path)
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []

        # 提取所有文本框（标题、正文等）
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
            # 提取幻灯片中的表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cells))

        if slide_parts:
            parts.append(f"--- Slide {i} ---")
            parts.extend(slide_parts)

    if not parts:
        raise ValueError(f"无法从 PPTX 文件中提取文本: {path.name}")

    return ParsedDocument(
        text="\n\n".join(parts),
        metadata={"slide_count": str(len(prs.slides))},
        source_format="pptx",
    )


# ── Markdown 解析 ─────────────────────────────────────────


def parse_markdown(path: Path) -> ParsedDocument:
    """解析 Markdown 文件，自动提取 YAML frontmatter 元数据。

    YAML frontmatter 示例（文件开头的 --- 块）:
        ---
        title: 我的文档
        author: 张三
        tags: python, rag
        ---

        正文内容从这里开始...

    frontmatter 中的内容会被提取到 metadata 字段，
    对检索和分类有帮助。
    """
    encoding = _detect_encoding(path)
    raw = path.read_bytes().decode(encoding, errors="replace")
    text = raw.strip()

    if not text:
        raise ValueError(f"Markdown 文件为空: {path.name}")

    metadata, body = _extract_frontmatter(text)

    return ParsedDocument(
        text=body.strip(),
        metadata={**metadata, "encoding": encoding},
        source_format="markdown",
    )


# 匹配开头 --- ... --- 块的简单正则（不依赖 yaml 库，更轻量）
_FRONTMATTER_RE = re.compile(r"^---\s*\n(.*?)\n---\s*\n", re.DOTALL)


def _extract_frontmatter(text: str) -> tuple[dict[str, str], str]:
    """从 Markdown 文本中提取 YAML frontmatter。

    返回:
        (metadata, body) —— metadata 是键值对字典，body 是去掉 frontmatter 后的正文。
    """
    match = _FRONTMATTER_RE.match(text)
    if not match:
        return {}, text

    fm_block = match.group(1)
    body = text[match.end():]

    metadata: dict[str, str] = {}
    current_key: str | None = None

    for line in fm_block.split("\n"):
        # 匹配 "key: value" 格式
        kv_match = re.match(r"^(\w[\w\s_-]*?):\s*(.*)", line)
        if kv_match:
            key = kv_match.group(1).strip()
            value = kv_match.group(2).strip().strip("\"'")
            if key:
                metadata[key] = value
                current_key = key
        elif current_key and line.strip().startswith("- "):
            # 列表项追加到上一个 key（用逗号拼接）
            item = line.strip()[2:].strip().strip("\"'")
            existing = metadata.get(current_key, "")
            metadata[current_key] = f"{existing}, {item}" if existing else item

    return metadata, body


# ── 纯文本解析 ────────────────────────────────────────────


def parse_text_file(path: Path) -> ParsedDocument:
    """解析纯文本文件（.txt），自动探测编码。"""
    encoding = _detect_encoding(path)
    text = path.read_text(encoding=encoding)
    cleaned = text.strip()

    if not cleaned:
        raise ValueError(f"文本文件为空: {path.name}")

    return ParsedDocument(
        text=cleaned,
        metadata={"encoding": encoding},
        source_format="text",
    )


# ── 统一入口 ──────────────────────────────────────────────


def parse_file(path: Path) -> ParsedDocument:
    """根据文件扩展名自动选择合适的解析器。

    这是外部模块调用的唯一入口，内部根据后缀分发到具体解析函数。

    Raises:
        ValueError: 不支持的文件类型或解析失败
    """
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix == ".doc":
        return parse_doc(path)
    if suffix == ".docx":
        return parse_docx(path)
    if suffix in {".xlsx", ".xls"}:
        return parse_xlsx(path)
    if suffix in {".pptx", ".ppt"}:
        return parse_pptx(path)
    if suffix in {".md", ".markdown"}:
        return parse_markdown(path)
    if suffix == ".txt":
        return parse_text_file(path)

    raise ValueError(f"不支持的文件类型: {suffix}")


# ── 文件收集 ──────────────────────────────────────────────


def collect_files(directory: Path, recursive: bool = True,
                  exclude_dirs: list[str] | None = None) -> list[Path]:
    """扫描目录，返回所有支持文件的排序列表。

    参数:
        directory: 要扫描的目录
        recursive: 是否递归扫描子目录
        exclude_dirs: 要跳过的目录名列表（路径中任一部分匹配即跳过）

    返回:
        按路径排序的文件列表（仅含 SUPPORTED_EXTENSIONS 中的类型）
    """
    if not directory.is_dir():
        raise NotADirectoryError(f"目录不存在: {directory}")

    exclude = exclude_dirs or []
    pattern = "**/*" if recursive else "*"
    files = [
        p
        for p in directory.glob(pattern)
        if p.is_file()
        and p.suffix.lower() in SUPPORTED_EXTENSIONS
        and not any(e in p.parts for e in exclude)
    ]
    return sorted(files)
